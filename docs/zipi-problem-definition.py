"""Zipi Problem Definition PPT Generator"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# === Design System ===
BG_BLACK = RGBColor(0x0A, 0x0A, 0x0A)
BG_DARK = RGBColor(0x12, 0x12, 0x12)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_300 = RGBColor(0xD1, 0xD5, 0xDB)
GRAY_400 = RGBColor(0x9C, 0xA3, 0xAF)
GRAY_500 = RGBColor(0x6B, 0x72, 0x80)
ACCENT_BLUE = RGBColor(0x60, 0xA5, 0xFA)
ACCENT_RED = RGBColor(0xF8, 0x71, 0x71)
ACCENT_GREEN = RGBColor(0x4A, 0xDE, 0x80)
ACCENT_AMBER = RGBColor(0xFB, 0xBF, 0x24)
ACCENT_PURPLE = RGBColor(0xA7, 0x8B, 0xFA)

FONT_TITLE = "Pretendard"
FONT_BODY = "Pretendard"
FONT_EN = "Inter"

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

def set_slide_bg(slide, color=BG_BLACK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name=FONT_BODY, line_spacing=1.5):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return tf

def add_multiline(slide, left, top, width, height, lines, font_size=18, color=WHITE, line_spacing=1.5, font_name=FONT_BODY):
    """lines: list of (text, color, bold, font_size_override)"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        if isinstance(line_data, str):
            text, c, b, fs = line_data, color, False, font_size
        else:
            text = line_data[0]
            c = line_data[1] if len(line_data) > 1 else color
            b = line_data[2] if len(line_data) > 2 else False
            fs = line_data[3] if len(line_data) > 3 else font_size
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(fs)
        p.font.color.rgb = c
        p.font.bold = b
        p.font.name = font_name
        p.space_after = Pt(4)
        if line_spacing != 1.0:
            p.line_spacing = Pt(fs * line_spacing)
    return tf

def add_accent_line(slide, left, top, width, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_number(slide, left, top, number, color=ACCENT_BLUE):
    add_text_box(slide, left, top, 1, 0.8, number, font_size=48, color=color, bold=True, font_name=FONT_EN)

# ============================
# SLIDE 1: Title
# ============================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_accent_line(slide, 1.5, 2.8, 2, ACCENT_BLUE)
add_text_box(slide, 1.5, 3.1, 13, 1.5, "모두가 만들 수 있는 시대,", font_size=44, color=WHITE, bold=True)
add_text_box(slide, 1.5, 4.2, 13, 1.5, "파는 사람의 기억력이 병목이 된다.", font_size=44, color=ACCENT_BLUE, bold=True)
add_text_box(slide, 1.5, 5.8, 10, 0.8, "Zipi  |  知彼知己 百戰百勝  |  Problem Definition", font_size=20, color=GRAY_400)

# ============================
# SLIDE 2: 시대 전환
# ============================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 1.5, 0.8, 10, 0.6, "01", font_size=14, color=ACCENT_BLUE, bold=True, font_name=FONT_EN)
add_text_box(slide, 1.5, 1.1, 13, 1, "시대 전환 — 모두가 빌더가 된다", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, 1.5, 2.0, 1.5, ACCENT_BLUE)

add_multiline(slide, 1.5, 2.5, 6, 5, [
    ("Andrej Karpathy:", ACCENT_BLUE, True, 20),
    ("제품 완성도 90% → 99% → 99.9%로 고도화", WHITE, False, 20),
    ("", WHITE, False, 10),
    ("AI 에이전트가 만들고, 코드를 짜고, 배포한다.", GRAY_300, False, 20),
    ("빌딩은 더 이상 병목이 아니다.", WHITE, True, 22),
])

add_multiline(slide, 8.5, 2.5, 6.5, 5, [
    ("그런데 우리가 간과한 것:", ACCENT_RED, True, 20),
    ("", WHITE, False, 10),
    ("좋은 제품은 저절로 팔리지 않는다.", WHITE, False, 20),
    ("파는 법은 제품력만큼 중요하다.", WHITE, True, 22),
    ("", WHITE, False, 10),
    ("그래서 모두가 세일즈맨이 되어야 한다.", GRAY_300, False, 20),
])

# ============================
# SLIDE 3: 왜 하필 세일즈인가 (점프 보완)
# ============================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 1.5, 0.8, 10, 0.6, "02", font_size=14, color=ACCENT_BLUE, bold=True, font_name=FONT_EN)
add_text_box(slide, 1.5, 1.1, 13, 1, "왜 하필 세일즈가 병목인가?", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, 1.5, 2.0, 1.5, ACCENT_BLUE)

add_multiline(slide, 1.5, 2.5, 6, 5.5, [
    ("빌딩 민주화 이후, 경쟁은 치열해진다.", GRAY_300, False, 20),
    ("디자인, 브랜딩, 유통 모두 병목 후보다.", GRAY_300, False, 20),
    ("", WHITE, False, 10),
    ("하지만 B2B 복합제품 영역에서는:", ACCENT_BLUE, True, 20),
    ("", WHITE, False, 10),
    ("  의사결정자가 여러 명이고,", WHITE, False, 20),
    ("  요구사항이 복잡하게 얽혀 있고,", WHITE, False, 20),
    ("  협상 과정에서 맥락이 누적된다.", WHITE, False, 20),
    ("", WHITE, False, 10),
    ("이 영역에서 PLG나 마케팅은 한계가 있다.", GRAY_400, False, 18),
    ("결국 사람이 만나서, 기억하고, 설득해야 한다.", WHITE, True, 22),
])

add_multiline(slide, 8.5, 2.5, 6.5, 5.5, [
    ("마케팅 콘텐츠도 AI가 만든다면,", GRAY_300, False, 20),
    ("신뢰는 어디서 오는가?", WHITE, True, 22),
    ("", WHITE, False, 10),
    ("리뷰 플랫폼? 조작 가능하다.", GRAY_400, False, 18),
    ("커뮤니티? 시간이 오래 걸린다.", GRAY_400, False, 18),
    ("PLG? 복합제품엔 한계가 있다.", GRAY_400, False, 18),
    ("", WHITE, False, 10),
    ("고가 + 복합 + 다수 의사결정자", ACCENT_AMBER, True, 20),
    ("= Physical Interaction이 대체 불가능", ACCENT_AMBER, True, 20),
])

# ============================
# SLIDE 4: 현장 경험
# ============================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 1.5, 0.8, 10, 0.6, "03", font_size=14, color=ACCENT_BLUE, bold=True, font_name=FONT_EN)
add_text_box(slide, 1.5, 1.1, 13, 1, "실제 경험 — Physical AI를 파는 현장", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, 1.5, 2.0, 1.5, ACCENT_BLUE)

# Left: 현장 상황
add_multiline(slide, 1.5, 2.5, 6, 6, [
    ("고객 미팅, 전화에서", WHITE, True, 22),
    ("말 한 마디 잘못하면 수천만 원이 오간다.", ACCENT_RED, True, 24),
    ("", WHITE, False, 10),
    ("팀장님: \"시뮬레이터 필요 없습니다.\"", GRAY_300, False, 20),
    ("", WHITE, False, 8),
    ("    ↓  CFO에게 보고", GRAY_500, False, 16),
    ("", WHITE, False, 8),
    ("CFO: \"눈에 보이는 산출물 가져와.\"", ACCENT_AMBER, False, 20),
    ("", WHITE, False, 8),
    ("    ↓  다시 우리에게", GRAY_500, False, 16),
    ("", WHITE, False, 8),
    ("팀장님: \"주신다는 얘기로 이해했는데,", ACCENT_RED, False, 20),
    ("          맞으시죠?\"", ACCENT_RED, False, 20),
])

# Right: 복잡성
add_multiline(slide, 8.5, 2.5, 6.5, 6, [
    ("우리 제품은 복잡하다.", WHITE, True, 22),
    ("", WHITE, False, 10),
    ("시뮬레이터를 줄 것인가?", GRAY_300, False, 20),
    ("어느 수준까지 정교화할 것인가?", GRAY_300, False, 20),
    ("소스코드는 포함하는가?", GRAY_300, False, 20),
    ("패키징과 이관은 별도 옵션인가?", GRAY_300, False, 20),
    ("", WHITE, False, 10),
    ("이 모든 맥락이 회의마다 쌓이고,", GRAY_400, False, 18),
    ("의사결정자마다 다르게 해석된다.", GRAY_400, False, 18),
])

# ============================
# SLIDE 5: 그 단 한 순간 (인과 체인 보완)
# ============================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 1.5, 0.8, 10, 0.6, "04", font_size=14, color=ACCENT_RED, bold=True, font_name=FONT_EN)
add_text_box(slide, 1.5, 1.1, 13, 1, "그 단 한 순간", font_size=40, color=ACCENT_RED, bold=True)
add_accent_line(slide, 1.5, 2.0, 1.5, ACCENT_RED)

# Left: 인과 체인
add_multiline(slide, 1.5, 2.5, 6, 6, [
    ("\"저번에 내가 뭐라고 했더라...?\"", WHITE, True, 28),
    ("", WHITE, False, 12),
    ("기억이 안 난다", GRAY_300, False, 20),
    ("  → 모호한 답변을 한다", GRAY_300, False, 20),
    ("  → 고객이 유리하게 해석한다", GRAY_300, False, 20),
    ("  → 나중에 뒤집기 어려워진다", GRAY_300, False, 20),
    ("  → 수천만 원의 팀 노력이", ACCENT_RED, True, 22),
    ("     공짜로 넘어간다", ACCENT_RED, True, 22),
])

# Right: 진짜 페인포인트
add_multiline(slide, 8.5, 2.5, 6.5, 6, [
    ("고객도 고마워하지 않는다.", WHITE, False, 20),
    ("우리 팀도 정당한 대가를 받지 못한다.", WHITE, False, 20),
    ("양쪽 모두 불행해지는 구조.", ACCENT_RED, True, 22),
    ("", WHITE, False, 12),
    ("세일즈맨의 진짜 페인:", ACCENT_AMBER, True, 22),
    ("", WHITE, False, 8),
    ("밤새 준비한 팀원들의 리소스를", WHITE, False, 20),
    ("내 말 한 마디 실수로 정당하게", WHITE, False, 20),
    ("보상받지 못하게 만든다는 죄책감.", ACCENT_RED, True, 22),
    ("", WHITE, False, 10),
    ("이건 단순한 매출 손실이 아니다.", GRAY_400, False, 18),
    ("동료에 대한 책임감의 문제다.", WHITE, True, 20),
])

# ============================
# SLIDE 6: 왜 이런 일이 생기는가 + 기존 도구 한계 (점프 보완)
# ============================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 1.5, 0.8, 10, 0.6, "05", font_size=14, color=ACCENT_BLUE, bold=True, font_name=FONT_EN)
add_text_box(slide, 1.5, 1.1, 13, 1, "왜 이런 일이 생기는가?", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, 1.5, 2.0, 1.5, ACCENT_BLUE)

add_multiline(slide, 1.5, 2.5, 6, 5.5, [
    ("세일즈맨의 기억력과 학습력이 한계에 부딪힌다.", WHITE, True, 22),
    ("", WHITE, False, 10),
    ("데이터는 이미 다 있다:", GRAY_300, False, 20),
    ("  녹음 파일, 회의록, 메모, 통화 기록...", GRAY_400, False, 18),
    ("", WHITE, False, 10),
    ("하지만 10초 안에 꺼내 쓸 수 없다.", ACCENT_BLUE, True, 22),
])

add_multiline(slide, 8.5, 2.5, 6.5, 5.5, [
    ("기존 도구들은 왜 안 되는가?", ACCENT_AMBER, True, 20),
    ("", WHITE, False, 10),
    ("Salesforce/CRM", WHITE, True, 18),
    ("  기록은 하지만, 맥락을 10초에 꺼내주지 못한다.", GRAY_400, False, 17),
    ("  수동 입력 → 바쁜 세일즈맨은 안 쓴다.", GRAY_400, False, 17),
    ("", WHITE, False, 8),
    ("Gong / Chorus", WHITE, True, 18),
    ("  영어권 중심, 분석 결과를 소비하기 어렵다.", GRAY_400, False, 17),
    ("  팀 매니저 관점 도구, 현장 세일즈맨 관점이 아니다.", GRAY_400, False, 17),
    ("", WHITE, False, 8),
    ("노션 / 메모", WHITE, True, 18),
    ("  기록 → 잊혀짐. 검색해도 맥락이 안 보인다.", GRAY_400, False, 17),
])

# ============================
# SLIDE 7: 이건 나만의 문제가 아니다 (N=1 보완)
# ============================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 1.5, 0.8, 10, 0.6, "06", font_size=14, color=ACCENT_BLUE, bold=True, font_name=FONT_EN)
add_text_box(slide, 1.5, 1.1, 13, 1, "이건 나만의 문제가 아니다", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, 1.5, 2.0, 1.5, ACCENT_BLUE)

add_multiline(slide, 1.5, 2.5, 6, 5.5, [
    ("시장이 검증하고 있다:", ACCENT_GREEN, True, 20),
    ("", WHITE, False, 10),
    ("Gong.io — $7.2B valuation", WHITE, True, 20),
    ("  세일즈 대화 분석 플랫폼", GRAY_400, False, 17),
    ("", WHITE, False, 8),
    ("Chorus (ZoomInfo) — $575M 인수", WHITE, True, 20),
    ("  대화 인텔리전스", GRAY_400, False, 17),
    ("", WHITE, False, 8),
    ("Clari, People.ai, Outreach...", WHITE, True, 20),
    ("  세일즈 인텔리전스 시장 성장 중", GRAY_400, False, 17),
    ("", WHITE, False, 10),
    ("이 시장은 이미 수조 원 규모다.", ACCENT_GREEN, True, 22),
])

add_multiline(slide, 8.5, 2.5, 6.5, 5.5, [
    ("하지만 아직 비어 있는 영역:", ACCENT_AMBER, True, 20),
    ("", WHITE, False, 10),
    ("기존 도구는 매니저 관점이다.", WHITE, False, 20),
    ("  \"우리 팀 세일즈 통화를 분석해줘\"", GRAY_400, False, 17),
    ("", WHITE, False, 8),
    ("Zipi는 현장 세일즈맨 관점이다.", ACCENT_BLUE, True, 20),
    ("  \"지금 고객이 전화했는데,", WHITE, False, 20),
    ("   10초 안에 맥락 알려줘\"", WHITE, True, 20),
    ("", WHITE, False, 10),
    ("AI Native 시대에는", GRAY_300, False, 20),
    ("이 기능이 개인의 기본 장비가 된다.", WHITE, True, 22),
])

# ============================
# SLIDE 8: 미래 — 커지는 시장
# ============================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 1.5, 0.8, 10, 0.6, "07", font_size=14, color=ACCENT_PURPLE, bold=True, font_name=FONT_EN)
add_text_box(slide, 1.5, 1.1, 13, 1, "Physical Interaction의 시대", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, 1.5, 2.0, 1.5, ACCENT_PURPLE)

add_multiline(slide, 1.5, 2.5, 13, 6, [
    ("모두가 빌더 → 모두가 세일즈맨이 되어야 한다", WHITE, True, 24),
    ("", WHITE, False, 10),
    ("AI가 만들고, AI가 마케팅하는 세상에서", GRAY_300, False, 20),
    ("B2B 복합제품의 신뢰는 결국 사람 간 대면에서 나온다.", WHITE, False, 20),
    ("", WHITE, False, 10),
    ("개인의 기억력 · 학습력이 병목이 된다.", ACCENT_PURPLE, True, 24),
    ("", WHITE, False, 10),
    ("맥락을 놓치는 것에 대한 답답함,", GRAY_300, False, 20),
    ("조직에서 무능력하다는 평가를 받는 것에 대한 고통.", GRAY_300, False, 20),
    ("이 페인은 AI 시대가 깊어질수록 계속 커진다.", WHITE, True, 22),
])

# ============================
# SLIDE 9: 솔루션 — Zipi
# ============================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 1.5, 0.8, 10, 0.6, "08", font_size=14, color=ACCENT_GREEN, bold=True, font_name=FONT_EN)
add_text_box(slide, 1.5, 1.1, 13, 1, "Zipi — 지피지기면 백전백승", font_size=40, color=WHITE, bold=True)
add_accent_line(slide, 1.5, 2.0, 1.5, ACCENT_GREEN)

add_text_box(slide, 1.5, 2.5, 13, 0.8, "고객을 탭 한 번 → 10초 안에 맥락 파악", font_size=28, color=ACCENT_GREEN, bold=True)

add_number(slide, 1.5, 3.6, "1", ACCENT_BLUE)
add_multiline(slide, 2.5, 3.7, 4, 1.5, [
    ("마지막 대화 요약", WHITE, True, 22),
    ("고객 대화 + 사내 회의 내용 통합", GRAY_400, False, 17),
])

add_number(slide, 6.5, 3.6, "2", ACCENT_BLUE)
add_multiline(slide, 7.5, 3.7, 4, 1.5, [
    ("예상 질문 + 추천 답변", WHITE, True, 22),
    ("다음 미팅에서 나올 질문을 미리 준비", GRAY_400, False, 17),
])

add_number(slide, 11.5, 3.6, "3", ACCENT_BLUE)
add_multiline(slide, 12.5, 3.7, 3, 1.5, [
    ("가격/조건 히스토리", WHITE, True, 22),
    ("누가 언제 뭐라고 했는지", GRAY_400, False, 17),
])

add_multiline(slide, 1.5, 5.8, 13, 2, [
    ("녹음 업로드 → 자동 분석 → 즉시 브리핑", GRAY_300, False, 20),
    ("", WHITE, False, 10),
    ("\"저번에 뭐라고 했더라?\" 를 없앤다.", WHITE, True, 28),
])

# ============================
# SLIDE 10: How We Built It — Ralph Loop
# ============================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 1.5, 0.8, 10, 0.6, "09", font_size=14, color=ACCENT_AMBER, bold=True, font_name=FONT_EN)
add_text_box(slide, 1.5, 1.1, 13, 1, "How We Built It — Ralph Loop", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, 1.5, 2.0, 1.5, ACCENT_AMBER)

# Left: image
pic = slide.shapes.add_picture(
    "C:/Project/Ralphbeforetest/pic.png",
    Inches(1.5), Inches(2.5), Inches(7), Inches(5.2)
)

# Right: story
add_multiline(slide, 9.2, 2.5, 5.8, 6, [
    ("가재 옷을 한 번도 안 입었습니다.", ACCENT_AMBER, True, 24),
    ("", WHITE, False, 10),
    ("해커톤 내내 Ralph Loop만 돌렸습니다.", WHITE, False, 20),
    ("AI 에이전트가 TDD로 코드를 짜고,", GRAY_300, False, 18),
    ("테스트하고, 커밋하고, 다음 마일스톤으로.", GRAY_300, False, 18),
    ("", WHITE, False, 10),
    ("20시간 돌리려고 했는데,", WHITE, False, 20),
    ("시간이 짧아서 아쉬웠습니다.", ACCENT_RED, True, 22),
    ("", WHITE, False, 10),
    ("이 앱은 그만큼 진심입니다.", WHITE, True, 24),
    ("제가 직접 겪은 문제이기 때문입니다.", GRAY_400, False, 18),
])

# ============================
# SLIDE 11: Closing
# ============================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_accent_line(slide, 1.5, 3.3, 2, ACCENT_BLUE)

add_text_box(slide, 1.5, 3.6, 13, 1.5,
    "세일즈의 그 단 한 순간을 지킨다.",
    font_size=44, color=WHITE, bold=True)

add_text_box(slide, 1.5, 5.3, 13, 1,
    "Zipi  |  知彼知己 百戰百勝",
    font_size=24, color=ACCENT_BLUE, bold=True)

# Save
output_path = "C:/Project/Ralphbeforetest/docs/zipi-problem-definition.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
